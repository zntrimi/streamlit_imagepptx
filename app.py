import streamlit as st
import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import math

def create_pptx_from_images(images):
    # 新しいプレゼンテーションを作成
    prs = Presentation()
    
    # スライドサイズを16:9に設定（2880×1620ピクセルに対応）
    prs.slide_width = Inches(13.3333)
    prs.slide_height = Inches(7.5)
    
    # 空白のスライドレイアウトを使用
    blank_slide_layout = prs.slide_layouts[6]
    
    # 各画像を新しいスライドに追加
    for uploaded_file in images:
        slide = prs.slides.add_slide(blank_slide_layout)
        # 画像をバイトストリームとして読み込む
        image_stream = io.BytesIO(uploaded_file.getvalue())
        slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
    
    # プレゼンテーションをバイトストリームとして保存
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    return pptx_stream.getvalue()

def show_image_grid(images):
    # 5列のグリッドで画像を表示
    num_images = len(images)
    num_rows = math.ceil(num_images / 5)
    
    for row in range(num_rows):
        cols = st.columns(5)
        for col_idx in range(5):
            img_idx = row * 5 + col_idx
            if img_idx < num_images:
                with cols[col_idx]:
                    st.image(images[img_idx], caption=f"Image {img_idx + 1}", use_column_width=True)

def main():
    st.title("Images to PowerPoint Converter")
    st.write("Upload your images and convert them into a PowerPoint presentation.")
    
    # 複数の画像をアップロード
    uploaded_files = st.file_uploader("Choose images", type=["png", "jpg", "jpeg"], accept_multiple_files=True)
    
    if uploaded_files:
        st.write(f"Number of images uploaded: {len(uploaded_files)}")
        
        # すべての画像をプレビュー表示
        st.subheader("Image Preview")
        show_image_grid(uploaded_files)
        
        if st.button("Convert to PowerPoint"):
            with st.spinner("Creating PowerPoint presentation..."):
                try:
                    pptx_bytes = create_pptx_from_images(uploaded_files)
                    
                    # ダウンロードボタンを表示
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_bytes,
                        file_name="presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    st.success("PowerPoint presentation created successfully!")
                except Exception as e:
                    st.error(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    main() 