import os
from pptx import Presentation
from pptx.util import Inches

# 1. 新しいプレゼンテーションを作成
prs = Presentation()

# 2. スライドサイズを16:9に設定（2880×1620ピクセルに対応）
prs.slide_width = Inches(13.3333)
prs.slide_height = Inches(7.5)

# 3. 画像ファイル名を1から100までの順に取得
image_files = ['{}.png'.format(i) for i in range(1, 101)]

# 4. 空白のスライドレイアウトを使用
blank_slide_layout = prs.slide_layouts[6]

# 5. 各画像を新しいスライドに追加
for image in image_files:
    if os.path.exists(image):
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(image, 0, 0, width=prs.slide_width, height=prs.slide_height)
    else:
        print('画像ファイルが見つかりません: {}'.format(image))

# 6. プレゼンテーションを保存
prs.save('output.pptx')